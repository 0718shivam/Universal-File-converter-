from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from PIL import Image, ImageDraw, ImageFont
import io
import os
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
import fitz  # PyMuPDF
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4, legal
from reportlab.lib.utils import ImageReader
import tempfile
import zipfile
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import pandas as pd
from openpyxl import load_workbook
import pytesseract

def _add_image_to_docx(doc, image_bytes, width=None, height=None):
    # Helper to add images to docx, resizing if necessary
    try:
        image_stream = io.BytesIO(image_bytes)
        # Attempt to open as PIL Image to get dimensions
        pil_img = Image.open(image_stream)
        img_width, img_height = pil_img.size
        image_stream.seek(0) # Reset stream for docx
        
        # Calculate aspect ratio
        aspect_ratio = img_width / img_height

        if width and not height:
            height = width / aspect_ratio
        elif height and not width:
            width = height * aspect_ratio
        elif not width and not height:
            # Default size if neither width nor height is specified
            width = Inches(6)
            height = width / aspect_ratio
        
        doc.add_picture(image_stream, width=width, height=height)
    except Exception as e:
        print(f"Error adding image to docx: {e}")

def pdf_to_images(pdf_file, output_format='PNG', dpi=200):
    """Convert PDF to images using PyMuPDF"""
    try:
        pdf_bytes = pdf_file.read()
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")

        output_files = []

        # Calculate zoom factor based on DPI
        zoom = dpi / 72  # 72 is the default DPI
        mat = fitz.Matrix(zoom, zoom)

        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            pix = page.get_pixmap(matrix=mat)

            # Convert pixmap to PIL Image
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            output = io.BytesIO()

            # Convert RGBA to RGB if saving as JPEG
            if output_format.upper() in ['JPEG', 'JPG'] and img.mode in ('RGBA', 'LA', 'P'):
                rgb_img = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                if img.mode in ('RGBA', 'LA'):
                    rgb_img.paste(img, mask=img.split()[-1])
                else:
                    rgb_img.paste(img)
                img = rgb_img

            img.save(output, format=output_format, quality=95)
            output.seek(0)
            output_files.append({
                'data': output.getvalue(),
                'filename': f'page_{page_num + 1}.{output_format.lower()}'
            })

        pdf_document.close()
        return output_files
    except Exception as e:
        raise Exception(f"Error converting PDF to images: {str(e)}")

def images_to_pdf(image_files, page_size='A4'):
    """Convert multiple images to a single PDF"""
    try:
        output = io.BytesIO()

        # Set page size
        if page_size.upper() == 'A4':
            size = A4
        elif page_size.upper() == 'LETTER':
            size = letter
        elif page_size.upper() == 'LEGAL':
            size = legal
        else:
            size = A4

        c = canvas.Canvas(output, pagesize=size)
        page_width, page_height = size

        for image_file in image_files:
            try:
                img = Image.open(io.BytesIO(image_file.read()))

                # Calculate scaling to fit page while maintaining aspect ratio
                img_width, img_height = img.size
                width_ratio = page_width / img_width
                height_ratio = page_height / img_height
                scale_ratio = min(width_ratio, height_ratio)

                new_width = img_width * scale_ratio
                new_height = img_height * scale_ratio

                # Center image on page
                x = (page_width - new_width) / 2
                y = (page_height - new_height) / 2

                # Save image to temporary BytesIO
                img_buffer = io.BytesIO()
                img.save(img_buffer, format='PNG')
                img_buffer.seek(0)

                # Draw image on canvas
                c.drawImage(ImageReader(img_buffer), x, y, new_width, new_height)
                c.showPage()

            except Exception as e:
                print(f"Error processing image: {str(e)}")
                continue

        c.save()
        output.seek(0)
        return output
    except Exception as e:
        raise Exception(f"Error converting images to PDF: {str(e)}")

def merge_pdfs(pdf_files):
    """Merge multiple PDFs into one"""
    try:
        merger = PdfMerger()

        for pdf_file in pdf_files:
            pdf_bytes = io.BytesIO(pdf_file.read())
            merger.append(pdf_bytes)

        output = io.BytesIO()
        merger.write(output)
        merger.close()
        output.seek(0)

        return output
    except Exception as e:
        raise Exception(f"Error merging PDFs: {str(e)}")

def split_pdf(pdf_file, split_type='all', page_range=None):
    """Split PDF into individual pages or by range"""
    try:
        pdf_bytes = io.BytesIO(pdf_file.read())
        reader = PdfReader(pdf_bytes)
        total_pages = len(reader.pages)

        output_files = []

        if split_type == 'all':
            # Split into individual pages
            for i in range(total_pages):
                writer = PdfWriter()
                writer.add_page(reader.pages[i])

                output = io.BytesIO()
                writer.write(output)
                output.seek(0)

                output_files.append({
                    'data': output.getvalue(),
                    'filename': f'page_{i+1}.pdf'
                })

        elif split_type == 'range' and page_range:
            # Split by page range (e.g., "1-3,5,7-9")
            ranges = page_range.split(',')
            for idx, r in enumerate(ranges):
                writer = PdfWriter()

                if '-' in r:
                    start, end = map(int, r.split('-'))
                    for i in range(start-1, min(end, total_pages)):
                        writer.add_page(reader.pages[i])
                else:
                    page_num = int(r) - 1
                    if 0 <= page_num < total_pages:
                        writer.add_page(reader.pages[page_num])

                output = io.BytesIO()
                writer.write(output)
                output.seek(0)

                output_files.append({
                    'data': output.getvalue(),
                    'filename': f'split_{idx+1}.pdf'
                })

        return output_files
    except Exception as e:
        raise Exception(f"Error splitting PDF: {str(e)}")

def compress_pdf(pdf_file, quality=60, dpi=150):
    try:
        src_bytes = pdf_file.read()
        doc = fitz.open(stream=src_bytes, filetype="pdf")
        has_images = False
        for i in range(doc.page_count):
            if doc.get_page_images(i):
                has_images = True
                break

        if not has_images:
            out = io.BytesIO()
            doc.save(out, garbage=4, deflate=True)
            doc.close()
            out.seek(0)
            return out

        out_doc = fitz.open()
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        for i in range(doc.page_count):
            page = doc[i]
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=int(quality), optimize=True)
            buf.seek(0)
            new_page = out_doc.new_page(width=page.rect.width, height=page.rect.height)
            new_page.insert_image(new_page.rect, stream=buf.getvalue())
        out = io.BytesIO()
        out_doc.save(out, garbage=4, deflate=True)
        out_doc.close()
        doc.close()
        out.seek(0)
        return out
    except Exception as e:
        raise Exception(f"Error compressing PDF: {str(e)}")

def extract_images(pdf_file, output_format='PNG'):
    """Extract embedded images from a PDF and return as list of bytes"""
    try:
        pdf_bytes = pdf_file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        images = []

        for page_num in range(doc.page_count):
            for img in doc.get_page_images(page_num):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                img_ext = base_image.get("ext", output_format.lower())

                # Convert to desired format if needed
                pil_img = Image.open(io.BytesIO(image_bytes))
                out_buf = io.BytesIO()
                save_format = output_format.upper()
                if save_format in ["JPEG", "JPG"] and pil_img.mode in ("RGBA", "LA", "P"):
                    rgb_img = Image.new('RGB', pil_img.size, (255, 255, 255))
                    if pil_img.mode == 'P':
                        pil_img = pil_img.convert('RGBA')
                    rgb_img.paste(pil_img, mask=pil_img.split()[-1] if pil_img.mode in ('RGBA', 'LA') else None)
                    pil_img = rgb_img
                pil_img.save(out_buf, format='JPEG' if save_format == 'JPG' else save_format)
                out_buf.seek(0)

                images.append({
                    'data': out_buf.getvalue(),
                    'filename': f'page_{page_num + 1}_img_{xref}.{("jpg" if save_format in ["JPEG", "JPG"] else save_format.lower())}'
                })

        doc.close()
        return images
    except Exception as e:
        raise Exception(f"Error extracting images: {str(e)}")

def get_pdf_info(pdf_file):
    """Get PDF metadata and information"""
    try:
        pdf_bytes = io.BytesIO(pdf_file.read())
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")

        info = {
            'pages': pdf_document.page_count,
            'metadata': {}
        }

        metadata = pdf_document.metadata
        if metadata:
            info['metadata'] = {
                'title': metadata.get('title', 'N/A'),
                'author': metadata.get('author', 'N/A'),
                'subject': metadata.get('subject', 'N/A'),
                'creator': metadata.get('creator', 'N/A'),
                'producer': metadata.get('producer', 'N/A'),
            }

        # Get page sizes
        page_sizes = []
        for page_num in range(min(5, pdf_document.page_count)):
            page = pdf_document[page_num]
            rect = page.rect
            page_sizes.append({
                'width': float(rect.width),
                'height': float(rect.height)
            })
        info['page_sizes'] = page_sizes

        pdf_document.close()
        return info
    except Exception as e:
        raise Exception(f"Error getting PDF info: {str(e)}")

def delete_pdf_pages(pdf_file, pages_to_delete):
    """Delete specific pages from PDF
    pages_to_delete: list of page numbers (1-indexed) or string like "1,3,5" or "1-3,5"
    """
    try:
        pdf_bytes = io.BytesIO(pdf_file.read())
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        total_pages = pdf_document.page_count

        # Parse pages to delete
        pages_set = set()
        if isinstance(pages_to_delete, str):
            ranges = pages_to_delete.split(',')
            for r in ranges:
                r = r.strip()
                if '-' in r:
                    start, end = map(int, r.split('-'))
                    pages_set.update(range(start, end + 1))
                else:
                    pages_set.add(int(r))
        elif isinstance(pages_to_delete, list):
            pages_set = set(pages_to_delete)

        # Delete pages in reverse order to avoid index shifting
        pages_to_delete_list = sorted(pages_set, reverse=True)
        for page_num in pages_to_delete_list:
            if 1 <= page_num <= total_pages:
                pdf_document.delete_page(page_num - 1)  # Convert to 0-indexed

        output = io.BytesIO()
        pdf_document.save(output)
        pdf_document.close()
        output.seek(0)

        return output
    except Exception as e:
        raise Exception(f"Error deleting PDF pages: {str(e)}")

def pdf_to_word(pdf_file):
    """Convert PDF to Word (DOCX) using PyMuPDF and python-docx"""
    try:
        doc = Document()
        pdf_bytes = pdf_file.read()
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")

        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]

            # Extract text
            text = page.get_text()
            if text:
                doc.add_paragraph(text)
            
            # Extract images (basic approach)
            for img_index, img in enumerate(pdf_document.get_page_images(page_num)):
                xref = img[0] # xref of the image
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]

                # Try to add image to document
                _add_image_to_docx(doc, image_bytes, width=Inches(6))

        pdf_document.close()

        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        return output
    except Exception as e:
        raise Exception(f"Error converting PDF to Word: {str(e)}")

def pdf_to_text(pdf_file):
    """Extract all text from a PDF file"""
    try:
        pdf_bytes = pdf_file.read()
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        text_content = []
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            text_content.append(page.get_text())
            
        pdf_document.close()
        return "\n".join(text_content)
    except Exception as e:
        raise Exception(f"Error extracting text from PDF: {str(e)}")

def pdf_to_excel(pdf_file):
    """Convert PDF text content to an Excel (XLSX) file"""
    try:
        pdf_bytes = pdf_file.read()
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")

        all_text = []
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            all_text.append(page.get_text())

        pdf_document.close()

        # Create a DataFrame from the extracted text
        df = pd.DataFrame(all_text, columns=['Extracted Text'])

        output = io.BytesIO()
        # Save to Excel, without index
        df.to_excel(output, index=False, engine='openpyxl')
        output.seek(0)
        return output
    except Exception as e:
        raise Exception(f"Error converting PDF to Excel: {str(e)}")

def excel_to_pdf(excel_file, page_size='A4'):
    """Convert an Excel (XLSX) file to PDF"""
    try:
        df = pd.read_excel(excel_file)

        output = io.BytesIO()
        
        # Set page size
        if page_size.upper() == 'A4':
            size = A4
        elif page_size.upper() == 'LETTER':
            size = letter
        elif page_size.upper() == 'LEGAL':
            size = legal
        else:
            size = A4

        c = canvas.Canvas(output, pagesize=size)
        width, height = size

        # Convert DataFrame to a string representation for PDF
        # This is a basic conversion; for complex layouts, more advanced rendering is needed.
        text_content = df.to_string(index=False)

        # Add text to PDF
        textobject = c.beginText()
        textobject.setTextOrigin(10, height - 20)
        textobject.setFont("Helvetica", 10)
        
        # Split text into lines to fit page width
        lines = text_content.split('\n')
        y_position = height - 20
        for line in lines:
            if y_position < 40:  # Check if close to bottom margin
                c.drawText(textobject) # Draw current text
                c.showPage() # Start new page
                textobject = c.beginText()
                textobject.setTextOrigin(10, height - 20)
                textobject.setFont("Helvetica", 10)
                y_position = height - 20
            textobject.textLine(line)
            y_position -= 12 # Line spacing

        c.drawText(textobject)
        c.save()
        output.seek(0)
        return output
    except Exception as e:
        raise Exception(f"Error converting Excel to PDF: {str(e)}")

def ocr_pdf(pdf_file):
    """Perform OCR on a PDF file to extract searchable text"""
    try:
        pdf_bytes = pdf_file.read()
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")

        full_text = []

        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            
            # Render page to image
            pix = page.get_pixmap(dpi=300)  # High DPI for better OCR accuracy
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            
            # Perform OCR using Tesseract
            text = pytesseract.image_to_string(img)
            full_text.append(text)

        pdf_document.close()
        return "\n".join(full_text)
    except Exception as e:
        raise Exception(f"Error performing OCR on PDF: {str(e)}")

def encrypt_pdf(pdf_file, password):
    try:
        src = io.BytesIO(pdf_file.read())
        reader = PdfReader(src)
        writer = PdfWriter()
        for p in reader.pages:
            writer.add_page(p)
        writer.encrypt(user_pwd=password, owner_pwd=password, use_128bit=True)
        out = io.BytesIO()
        writer.write(out)
        out.seek(0)
        return out
    except Exception as e:
        raise Exception(f"Error encrypting PDF: {str(e)}")

def decrypt_pdf(pdf_file, password):
    try:
        src = io.BytesIO(pdf_file.read())
        reader = PdfReader(src)
        if reader.is_encrypted:
            res = reader.decrypt(password)
            if res == 0:
                raise Exception("Incorrect password for PDF decryption.")
        else:
            src.seek(0)
            return src
        writer = PdfWriter()
        for p in reader.pages:
            writer.add_page(p)
        out = io.BytesIO()
        writer.write(out)
        out.seek(0)
        return out
    except Exception as e:
        raise Exception(f"Error decrypting PDF: {str(e)}")

def pdf_to_ppt(pdf_file, layout='blank'):
    """Convert PDF to PowerPoint presentation using PyMuPDF"""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        pdf_bytes = pdf_file.read()
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")

        # Create PowerPoint presentation
        prs = Presentation()

        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Convert PDF pages to images
        zoom = 200 / 72  # High DPI for quality
        mat = fitz.Matrix(zoom, zoom)

        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            pix = page.get_pixmap(matrix=mat)

            # Convert to PIL Image
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)

            # Convert PIL image to bytes
            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            # Calculate image size to fit slide while maintaining aspect ratio
            img_width, img_height = img.size
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # Calculate scaling
            width_ratio = slide_width / img_width
            height_ratio = slide_height / img_height
            scale_ratio = min(width_ratio, height_ratio) * 0.95  # 95% to add some padding

            new_width = int(img_width * scale_ratio)
            new_height = int(img_height * scale_ratio)

            # Center image
            left = (slide_width - new_width) / 2
            top = (slide_height - new_height) / 2

            # Add image to slide
            slide.shapes.add_picture(img_buffer, left, top, new_width, new_height)

        pdf_document.close()

        # Save presentation
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        return output
    except Exception as e:
        raise Exception(f"Error converting PDF to PPT: {str(e)}")

def rotate_pdf_pages(pdf_file, rotation=90, pages='all'):
    """Rotate PDF pages using PyMuPDF
    rotation: 90, 180, or 270 degrees clockwise
    pages: 'all' or list of page numbers (1-indexed) or string like "1,3,5"
    """
    try:
        pdf_bytes = io.BytesIO(pdf_file.read())
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        total_pages = pdf_document.page_count

        # Parse pages to rotate
        pages_set = set()
        if pages == 'all':
            pages_set = set(range(1, total_pages + 1))
        elif isinstance(pages, str):
            ranges = pages.split(',')
            for r in ranges:
                r = r.strip()
                if '-' in r:
                    start, end = map(int, r.split('-'))
                    pages_set.update(range(start, end + 1))
                else:
                    pages_set.add(int(r))
        elif isinstance(pages, list):
            pages_set = set(pages)

        # Rotate specified pages
        for page_num in range(total_pages):
            if (page_num + 1) in pages_set:
                page = pdf_document[page_num]
                page.set_rotation(rotation)

        output = io.BytesIO()
        pdf_document.save(output)
        pdf_document.close()
        output.seek(0)

        return output
    except Exception as e:
        raise Exception(f"Error rotating PDF pages: {str(e)}")

def add_watermark(pdf_file, text, opacity=0.2, font_size=36, pages='all'):
    try:
        pdf_stream = io.BytesIO(pdf_file.read())
        doc = fitz.open(stream=pdf_stream, filetype="pdf")
        total = doc.page_count
        sel = set()
        if pages == 'all':
            sel = set(range(1, total + 1))
        elif isinstance(pages, str):
            for r in pages.split(','):
                r = r.strip()
                if '-' in r:
                    s, e = map(int, r.split('-'))
                    sel.update(range(s, e + 1))
                else:
                    sel.add(int(r))
        elif isinstance(pages, list):
            sel = set(pages)
        for idx in range(total):
            if (idx + 1) in sel:
                page = doc[idx]
                rect = page.rect
                w = int(rect.width)
                h = int(rect.height)
                overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
                txt = Image.new("RGBA", (w, h), (0, 0, 0, 0))
                draw = ImageDraw.Draw(txt)
                try:
                    font = ImageFont.truetype("arial.ttf", font_size)
                except Exception:
                    font = ImageFont.load_default()
                draw.text((w // 2, h // 2), text, font=font, fill=(0, 0, 0, int(max(0, min(1, opacity)) * 255)), anchor="mm")
                txt = txt.rotate(45, resample=Image.BICUBIC)
                overlay = Image.alpha_composite(overlay, txt)
                buf = io.BytesIO()
                overlay.save(buf, format="PNG")
                buf.seek(0)
                page.insert_image(rect, stream=buf.getvalue())
        out = io.BytesIO()
        doc.save(out)
        doc.close()
        out.seek(0)
        return out
    except Exception as e:
        raise Exception(f"Error adding watermark: {str(e)}")
